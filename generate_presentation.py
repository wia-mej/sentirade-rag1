"""
Présentation académique Sentirade RAG — Focalisée ML — version finale.
Format : 10×5.62 in (même que l'exemple fourni).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x1B, 0x2A)
CARD    = RGBColor(0x16, 0x2A, 0x40)
CARD2   = RGBColor(0x0F, 0x22, 0x35)
CYAN    = RGBColor(0x00, 0xD4, 0xFF)
GREEN   = RGBColor(0x00, 0xE5, 0x96)
ORANGE  = RGBColor(0xFF, 0x8C, 0x00)
RED     = RGBColor(0xFF, 0x45, 0x6A)
PURPLE  = RGBColor(0xA0, 0x55, 0xFF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0xB0, 0xC4, 0xDE)
DARK    = RGBColor(0x0A, 0x12, 0x1E)

SW = Inches(10)
SH = Inches(5.62)

FOOTER = "ISITD 2024 — S4 Machine Learning   |   Avril 2026"

# ── Low-level helpers ──────────────────────────────────────────────────────────
def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, radius=False):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    return shp

def tb(slide, text, l, t, w, h,
       size=14, bold=False, italic=False,
       color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return txb

def header(slide, title, section_num=None, subtitle=None):
    """Barre de titre standard."""
    rect(slide, 0, 0, SW, Inches(0.9), CARD)
    rect(slide, 0, Inches(0.9), SW, Pt(3), CYAN)
    prefix = f"{section_num} — " if section_num else ""
    tb(slide, prefix + title,
       Inches(0.3), Inches(0.06), Inches(9.2), Inches(0.72),
       size=22, bold=True, color=CYAN)
    if subtitle:
        tb(slide, subtitle,
           Inches(0.3), Inches(0.68), Inches(9.4), Inches(0.28),
           size=10, italic=True, color=GREY)

def footer_bar(slide):
    rect(slide, 0, Inches(5.35), SW, Inches(0.27), DARK)
    tb(slide, FOOTER,
       Inches(0.3), Inches(5.36), Inches(9.4), Inches(0.26),
       size=8, color=GREY, align=PP_ALIGN.CENTER)

def bullet_list(slide, items, x, y, w, line_h=Inches(0.32)):
    """items: list of dicts {text, indent?, color?, size?, bold?}"""
    cy = y
    for item in items:
        indent = item.get("indent", 0)
        text   = item["text"]
        color  = item.get("color", WHITE)
        size   = item.get("size", 13)
        bold   = item.get("bold", False)
        prefix = "    " * indent + ("▸ " if indent else "• ")
        tb(slide, prefix + text, x, cy, w, line_h,
           size=size, color=color, bold=bold)
        cy += line_h + Inches(0.01 * (1 + indent))
    return cy

def stat_box(slide, l, t, w, h, value, label, color, sub=None):
    """Grande boîte de statistique centrée."""
    rect(slide, l, t, w, h, color)
    tb(slide, value, l, t, w, h * 0.55,
       size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    tb(slide, label, l, t + h * 0.5, w, h * 0.3,
       size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    if sub:
        tb(slide, sub, l, t + h * 0.78, w, h * 0.25,
           size=9, color=DARK, align=PP_ALIGN.CENTER, italic=True)

def formula_box(slide, l, t, w, h, text, color=CARD2):
    rect(slide, l, t, w, h, color)
    tb(slide, text, l + Inches(0.08), t + Inches(0.06),
       w - Inches(0.16), h - Inches(0.12),
       size=11, color=CYAN, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Titre
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
rect(s, 0, 0, SW, Inches(2.6), CARD)
rect(s, 0, Inches(2.6), SW, Pt(4), CYAN)

tb(s, "SENTIRADE", Inches(0.4), Inches(0.2), Inches(9.2), Inches(1.0),
   size=58, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
tb(s, "RAG", Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.85),
   size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Système de Trading Intelligent basé sur le RAG",
   Inches(0.4), Inches(1.95), Inches(9.2), Inches(0.5),
   size=15, italic=True, color=GREY, align=PP_ALIGN.CENTER)

badges = [
    ("Agent ReAct", CYAN), ("GMM", GREEN), ("XGBoost", ORANGE),
    ("ChromaDB", PURPLE), ("LLaMA 3.3 70B", RED), ("Electron", GREY),
]
bx = Inches(0.25)
for label, col in badges:
    rect(s, bx, Inches(2.82), Inches(1.53), Inches(0.38), col)
    tb(s, label, bx, Inches(2.82), Inches(1.53), Inches(0.38),
       size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    bx += Inches(1.59)

rect(s, 0, Inches(5.35), SW, Inches(0.27), DARK)
tb(s, "S1 · S2 · S3 · S4   |   ISITD 2024 — S4 Machine Learning   |   Avril 2026",
   Inches(0.3), Inches(5.36), Inches(9.4), Inches(0.26),
   size=8, color=GREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Plan
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "Plan de la présentation")
footer_bar(s)

sections = [
    ("01", "Problématique & Objectif",        "Pourquoi fusionner technique, news et ML ?"),
    ("02", "Les Données — Vue d'ensemble",     "Deux phases : simulées puis Kaggle réelles"),
    ("03", "Datasets Kaggle",                  "Stock Market Dataset + Financial Phrase Bank"),
    ("04", "Architecture RAG & Agent ReAct",   "ChromaDB, Groq LLaMA 3.3, boucle ReAct"),
    ("05", "Analyse Technique — Couche 1",     "RSI, Volatilité, MA Spread"),
    ("06", "Détection de Régime — GMM",        "Apprentissage non supervisé, 3 régimes"),
    ("07", "Sentiment RAG — Couche 2",         "Embeddings, retrieval, génération JSON"),
    ("08", "Feature Matrix & Fusion",          "Vecteur 7D hétérogène"),
    ("09", "XGBoost — Classifieur",            "Direction, AUC-ROC, split temporel"),
    ("10", "XGBoost — Régresseur",             "Amplitude, MAE, RMSE"),
    ("11", "Backtesting & Résultats",          "Sharpe Ratio, Win Rate, benchmark S&P 500"),
    ("12", "Pipeline complet",                 "De la donnée brute à la décision"),
    ("13", "Conclusion",                       ""),
]

col_w = Inches(4.7)
for i, (num, title, sub) in enumerate(sections):
    col = 0 if i < 6 else 1
    row = i % 6
    x = Inches(0.25) + col * Inches(5.0)
    y = Inches(1.05) + row * Inches(0.7)
    rect(s, x, y, col_w, Inches(0.6), CARD)
    tb(s, num, x + Inches(0.1), y + Inches(0.05), Inches(0.55), Inches(0.5),
       size=18, bold=True, color=CYAN)
    tb(s, title, x + Inches(0.7), y + Inches(0.04), Inches(3.8), Inches(0.32),
       size=12, bold=True, color=WHITE)
    if sub:
        tb(s, sub, x + Inches(0.7), y + Inches(0.32), Inches(3.8), Inches(0.24),
           size=9, color=GREY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Problématique
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "Problématique & Objectif", "01",
       subtitle="Comment prédire le comportement d'une action en combinant technique, actualités et IA ?")
footer_bar(s)

icons = [
    ("📊", "Analyse\ntechnique",    "RSI · Volatilité\nMA Spread",      CYAN),
    ("📰", "Actualités\nfinancières","RAG + ChromaDB\n+ Groq LLaMA",     GREEN),
    ("🤖", "Intelligence\nArtificielle","Agent ReAct\nXGBoost",           ORANGE),
    ("📈", "Prédiction\ndouble",    "Direction\n+ Amplitude",             PURPLE),
]
bx = Inches(0.3)
for icon, title, sub, col in icons:
    rect(s, bx, Inches(1.1), Inches(2.15), Inches(3.5), CARD)
    tb(s, icon, bx, Inches(1.15), Inches(2.15), Inches(0.7),
       size=30, align=PP_ALIGN.CENTER)
    tb(s, title, bx, Inches(1.85), Inches(2.15), Inches(0.65),
       size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
    tb(s, sub, bx, Inches(2.5), Inches(2.15), Inches(0.8),
       size=10, color=GREY, align=PP_ALIGN.CENTER)
    rect(s, bx, Inches(4.55), Inches(2.15), Pt(3), col)
    bx += Inches(2.35)

tb(s, "Enjeu : les marchés réagissent autant aux chiffres qu'aux mots — aucun modèle isolé ne capture les deux.",
   Inches(0.3), Inches(4.75), Inches(9.4), Inches(0.5),
   size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Les Données : Vue d'ensemble (2 phases)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "Les Données — Vue d'ensemble", "02",
       subtitle="Deux phases d'expérimentation : données simulées → données Kaggle réelles")
footer_bar(s)

# Phase 1
rect(s, Inches(0.2), Inches(1.02), Inches(4.65), Inches(2.0), CARD)
rect(s, Inches(0.2), Inches(1.02), Pt(5), Inches(2.0), ORANGE)
tb(s, "Phase 1 — Données Simulées", Inches(0.38), Inches(1.08), Inches(4.35), Inches(0.32),
   size=12, bold=True, color=ORANGE)
bullet_list(s, [
    {"text": "Problème : yfinance bloqué par le réseau universitaire", "color": GREY},
    {"text": "Solution : génération de données OHLCV synthétiques",   "color": WHITE},
    {"text": "Distribution statistique proche de la réalité",          "color": GREY},
    {"text": "AAPL · TSLA · NVDA — 1305 jours chacun",               "color": WHITE, "bold": True},
    {"text": "AUC-ROC obtenu : 0.379 (proche du hasard)",             "color": RED, "bold": True},
], Inches(0.38), Inches(1.43), Inches(4.35), line_h=Inches(0.27))

# Flèche de transition
tb(s, "⬇ Passage aux données réelles", Inches(0.2), Inches(3.1), Inches(4.65), Inches(0.32),
   size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# Phase 2
rect(s, Inches(0.2), Inches(3.45), Inches(4.65), Inches(1.7), CARD)
rect(s, Inches(0.2), Inches(3.45), Pt(5), Inches(1.7), GREEN)
tb(s, "Phase 2 — Données Kaggle Réelles", Inches(0.38), Inches(3.51), Inches(4.35), Inches(0.32),
   size=12, bold=True, color=GREEN)
bullet_list(s, [
    {"text": "Stock Market Dataset — OHLCV réels 2018–2020",          "color": WHITE, "bold": True},
    {"text": "Financial Phrase Bank — 4846 headlines annotées",        "color": WHITE, "bold": True},
    {"text": "1500 news nettoyées ingérées dans ChromaDB",             "color": GREY},
    {"text": "AUC-ROC : 0.379  →  0.645  ✓",                         "color": GREEN, "bold": True},
], Inches(0.38), Inches(3.88), Inches(4.35), line_h=Inches(0.27))

# Droite — stats globales
stats = [
    ("0.379\n→ 0.645", "⬆ AUC-ROC\naprès Kaggle",  GREEN),
    ("1500",           "News nettoyées\nChromaDB",   CYAN),
    ("566",            "Jours boursiers\npar ticker", ORANGE),
    ("150",            "Lignes feature\nmatrix finale",PURPLE),
]
sx = Inches(5.1)
sy = Inches(1.02)
for i, (val, lbl, col) in enumerate(stats):
    cx = sx + (i % 2) * Inches(2.3)
    cy = sy + (i // 2) * Inches(1.9)
    stat_box(s, cx, cy, Inches(2.05), Inches(1.65), val, lbl, col)

tb(s, "Split temporel strict : 70% train · 15% val · 15% test — jamais de shuffle aléatoire",
   Inches(0.2), Inches(5.1), Inches(9.6), Inches(0.25),
   size=9, italic=True, color=RED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Datasets Kaggle (nouveau)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "Datasets Kaggle — Détail", "03",
       subtitle="Deux datasets publics ont remplacé les données simulées pour améliorer la qualité du modèle")
footer_bar(s)

# Dataset 1 — Stock Market
rect(s, Inches(0.2), Inches(1.02), Inches(4.65), Inches(3.85), CARD)
rect(s, Inches(0.2), Inches(1.02), Pt(5), Inches(3.85), CYAN)
tb(s, "① Stock Market Dataset", Inches(0.38), Inches(1.08), Inches(4.3), Inches(0.32),
   size=13, bold=True, color=CYAN)
tb(s, "Prix OHLCV réels — 2018 à 2020",
   Inches(0.38), Inches(1.42), Inches(4.3), Inches(0.28),
   size=10, color=GREY, italic=True)
bullet_list(s, [
    {"text": "566 jours boursiers par ticker",                               "color": WHITE},
    {"text": "3 tickers d'entraînement : AAPL · TSLA · NVDA",               "color": WHITE, "bold": True},
    {"text": "Colonnes : Date · Open · High · Low · Close · Volume",          "color": GREY},
    {"text": "Remplace les données OHLCV synthétiques de la Phase 1",        "color": GREY},
    {"text": "Utilisé pour :", "bold": True, "color": WHITE},
    {"indent": 1, "text": "Calcul RSI, Volatilité, MA Spread",               "color": GREY, "size": 10},
    {"indent": 1, "text": "Entraînement GMM (régimes de marché)",            "color": GREY, "size": 10},
    {"indent": 1, "text": "Cible XGBoost : retour J+1 réel",                "color": GREY, "size": 10},
    {"indent": 1, "text": "Backtest sur prix de marché authentiques",        "color": GREY, "size": 10},
], Inches(0.38), Inches(1.73), Inches(4.3), line_h=Inches(0.27))

# Dataset 2 — Financial Phrase Bank
rect(s, Inches(5.15), Inches(1.02), Inches(4.65), Inches(3.85), CARD)
rect(s, Inches(5.15), Inches(1.02), Pt(5), Inches(3.85), GREEN)
tb(s, "② Financial Phrase Bank", Inches(5.33), Inches(1.08), Inches(4.3), Inches(0.32),
   size=13, bold=True, color=GREEN)
tb(s, "Kaggle — 4846 headlines financières annotées",
   Inches(5.33), Inches(1.42), Inches(4.3), Inches(0.28),
   size=10, color=GREY, italic=True)
bullet_list(s, [
    {"text": "4846 phrases annotées manuellement",                            "color": WHITE},
    {"text": "3 labels : positif · négatif · neutre",                        "color": WHITE, "bold": True},
    {"text": "Source : rapports financiers, presse économique",               "color": GREY},
    {"text": "Nettoyage & filtrage : 1500 headlines retenues",                "color": WHITE, "bold": True},
    {"text": "Ingérées dans ChromaDB (financial_news) :",                     "bold": True, "color": WHITE},
    {"indent": 1, "text": "Embedding : all-MiniLM-L6-v2 (384D)",             "color": GREY, "size": 10},
    {"indent": 1, "text": "Métadonnées : ticker · date · label · source",    "color": GREY, "size": 10},
    {"indent": 1, "text": "Enrichit le RAG avec des exemples annotés réels", "color": GREY, "size": 10},
    {"text": "Impact : AUC-ROC 0.379  →  0.645  (+70%)",                     "color": GREEN, "bold": True},
], Inches(5.33), Inches(1.73), Inches(4.3), line_h=Inches(0.27))

# Barre de comparaison AUC
rect(s, Inches(0.2), Inches(5.0), Inches(9.6), Inches(0.28), CARD2)
tb(s, "AUC-ROC avec données simulées : 0.379",
   Inches(0.35), Inches(5.02), Inches(4.4), Inches(0.24),
   size=9, color=RED, bold=True)
tb(s, "AUC-ROC avec données Kaggle : 0.645  ✓  (+70%)",
   Inches(5.0), Inches(5.02), Inches(4.6), Inches(0.24),
   size=9, color=GREEN, bold=True, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Architecture RAG & Agent ReAct
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "Architecture RAG & Agent ReAct", "04",
       subtitle="RAG = Retrieval-Augmented Generation — l'agent récupère des faits avant de générer")
footer_bar(s)

# Boucle ReAct — 3 itérations
iters = [
    ("Itération 1", "REASONING\nJe dois connaître\nle régime de marché",
     "ACTION :\nget_regime()\n→ GMM label", CYAN),
    ("Itération 2", "REASONING\nJe dois chercher\nles news pertinentes",
     "ACTION :\nsearch_news_db()\n→ ChromaDB top-N", GREEN),
    ("Itération 3", "REASONING\nJ'ai assez d'infos\npour générer un signal",
     "ACTION :\nGroq LLaMA\n→ JSON structuré", ORANGE),
]
bx = Inches(0.2)
for ititle, reason, action, col in iters:
    rect(s, bx, Inches(1.05), Inches(2.9), Inches(3.5), CARD)
    tb(s, ititle, bx + Inches(0.05), Inches(1.1), Inches(2.8), Inches(0.3),
       size=10, bold=True, color=col)
    rect(s, bx + Inches(0.08), Inches(1.42), Inches(1.3), Inches(1.35), CARD2)
    tb(s, reason, bx + Inches(0.1), Inches(1.45), Inches(1.25), Inches(1.3),
       size=9, color=WHITE)
    rect(s, bx + Inches(1.5), Inches(1.42), Inches(1.32), Inches(1.35), CARD2)
    tb(s, action, bx + Inches(1.52), Inches(1.45), Inches(1.28), Inches(1.3),
       size=9, color=col)
    bx += Inches(3.1)

# Composants clés (bas)
components = [
    ("ChromaDB",        "Base vectorielle locale\n1500+ docs, 384 dimensions",   CYAN),
    ("all-MiniLM-L6-v2","Modèle d'embedding\nTexte → vecteurs 384D",             GREEN),
    ("Groq LLaMA 3.3",  "Analyse sentiment\n→ JSON : signal + confiance",         ORANGE),
    ("decision_log",    "Traçabilité complète\n450+ lignes de log",               PURPLE),
]
bx = Inches(0.2)
for ctitle, csub, col in components:
    rect(s, bx, Inches(4.75), Inches(2.35), Inches(0.68), col)
    tb(s, ctitle, bx + Inches(0.06), Inches(4.77), Inches(2.2), Inches(0.28),
       size=11, bold=True, color=DARK)
    tb(s, csub, bx + Inches(0.06), Inches(5.04), Inches(2.2), Inches(0.36),
       size=8, color=DARK)
    bx += Inches(2.47)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Couche 1 : Analyse Technique
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "Analyse Technique — Couche 1", "05",
       subtitle="Extraction de 3 features quantitatives à partir des séries OHLCV")
footer_bar(s)

features = [
    ("RSI — Relative Strength Index (période 14)", CYAN,
     "Mesure la vitesse et l'amplitude des variations de prix.\n"
     "RSI < 30 → survente (signal d'achat potentiel)\n"
     "RSI > 70 → surachat (signal de vente potentiel)",
     "RSI = 100 - 100 / (1 + RS)\nRS = Gain moyen(14j) / Perte moyenne(14j)",
     "Implémentation manuelle (sans TA-Lib)\nCalcul des gains/pertes moyens sur fenêtre glissante"),

    ("Volatilité — Écart-type glissant (14j)", ORANGE,
     "Mesure le risque instantané du titre.\n"
     "Forte volatilité → incertitude élevée → pondération réduite\n"
     "Utilisée par le GMM comme proxy du risque",
     "σ(t) = std( Δprix(t-13) … Δprix(t) )\nΔprix = (Close(t) - Close(t-1)) / Close(t-1)",
     "Feature clé pour le clustering GMM\nCorrelée négativement avec le régime haussier"),

    ("MA Spread — Croisement de moyennes mobiles", GREEN,
     "Capture la divergence court terme / long terme.\n"
     "MA Spread > 0 → tendance haussière de court terme\n"
     "MA Spread < 0 → tendance baissière",
     "Spread = (MA10 - MA50) / MA50\nNormalisé par MA50 pour comparabilité inter-tickers",
     "Signal de momentum structurel\nBase du signal technique envoyé au prompt LLM"),
]

y = Inches(1.02)
for title, col, desc, formula, note in features:
    rect(s, Inches(0.2), y, Inches(9.6), Inches(1.35), CARD)
    rect(s, Inches(0.2), y, Pt(5), Inches(1.35), col)
    tb(s, title, Inches(0.38), y + Inches(0.06), Inches(4.1), Inches(0.3),
       size=12, bold=True, color=col)
    tb(s, desc, Inches(0.38), y + Inches(0.38), Inches(4.0), Inches(0.9),
       size=9.5, color=WHITE)
    formula_box(s, Inches(4.55), y + Inches(0.06), Inches(3.15), Inches(0.65), formula)
    tb(s, note, Inches(4.55), y + Inches(0.75), Inches(3.15), Inches(0.56),
       size=8.5, color=GREY, italic=True)
    y += Inches(1.48)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — GMM Détection de Régime
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "Détection de Régime — Gaussian Mixture Model", "06",
       subtitle="Apprentissage non supervisé : identifier le contexte de marché avant d'analyser les news")
footer_bar(s)

# Gauche — principe
rect(s, Inches(0.2), Inches(1.02), Inches(5.3), Inches(4.05), CARD)
tb(s, "Pourquoi le GMM ?", Inches(0.35), Inches(1.08), Inches(5.0), Inches(0.32),
   size=12, bold=True, color=CYAN)
bullet_list(s, [
    {"text": "K-Means impose des frontières dures — inadapté aux marchés financiers", "color": GREY},
    {"text": "Le GMM offre une appartenance probabiliste (soft assignment)", "color": WHITE, "bold": True},
    {"text": "3 composantes (K=3) : Haussier · Calme · Baissier", "color": WHITE, "bold": True},
    {"text": "Chaque point suit un mélange de 3 gaussiennes multivariées", "color": GREY},
], Inches(0.35), Inches(1.43), Inches(5.0))

tb(s, "Entrée & Normalisation", Inches(0.35), Inches(2.82), Inches(5.0), Inches(0.3),
   size=12, bold=True, color=ORANGE)
formula_box(s, Inches(0.35), Inches(3.15), Inches(5.0), Inches(0.55),
            "X = [RSI, Volatilité, MA_Spread]  →  StandardScaler  →  GMM.fit_predict(X)")

tb(s, "Propriétés clés", Inches(0.35), Inches(3.78), Inches(5.0), Inches(0.28),
   size=12, bold=True, color=GREEN)
bullet_list(s, [
    {"text": "Re-fit dynamique à chaque nouveau ticker détecté", "color": GREY},
    {"text": "Résultat injecté comme feature dans la feature matrix", "color": WHITE},
    {"text": "Conditionne le prompt LLM (contexte de marché)", "color": CYAN},
], Inches(0.35), Inches(4.1), Inches(5.0), line_h=Inches(0.28))

# Droite — 3 régimes
regimes = [
    ("Régime 0", "Haussier  🟢", "RSI élevé · Faible volatilité · MA Spread positif\nMomentum positif soutenu — contexte favorable au BUY", GREEN),
    ("Régime 1", "Calme  ⚪",    "RSI neutre · Volatilité modérée · Spread proche 0\nMarché sans tendance claire — signal HOLD probable",   GREY),
    ("Régime 2", "Baissier  🔴", "RSI faible · Forte volatilité · MA Spread négatif\nPression vendeuse — contexte favorable au SELL",       RED),
]
ry = Inches(1.02)
for rnum, rname, rdesc, rcol in regimes:
    rect(s, Inches(5.7), ry, Inches(4.1), Inches(1.26), CARD)
    rect(s, Inches(5.7), ry, Pt(5), Inches(1.26), rcol)
    tb(s, rnum, Inches(5.88), ry + Inches(0.06), Inches(3.8), Inches(0.28),
       size=9, color=GREY)
    tb(s, rname, Inches(5.88), ry + Inches(0.3), Inches(3.8), Inches(0.32),
       size=13, bold=True, color=rcol)
    tb(s, rdesc, Inches(5.88), ry + Inches(0.63), Inches(3.8), Inches(0.6),
       size=9, color=WHITE)
    ry += Inches(1.37)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Sentiment RAG & LLM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "Sentiment RAG — Couche 2", "07",
       subtitle="Retrieval-Augmented Generation : ancrer le LLM dans des faits financiers récents")
footer_bar(s)

# Étape R
rect(s, Inches(0.2), Inches(1.02), Inches(4.6), Inches(3.95), CARD)
tb(s, "① RETRIEVAL — ChromaDB", Inches(0.35), Inches(1.08), Inches(4.3), Inches(0.3),
   size=12, bold=True, color=GREEN)
bullet_list(s, [
    {"text": "Collection : financial_news",                   "color": GREY},
    {"text": "Modèle d'embedding : all-MiniLM-L6-v2 (384D)", "color": WHITE},
    {"text": "Ingestion : yfinance.Ticker(ticker).news",       "color": GREY},
    {"text": "Chaque article → vecteur + métadonnées",        "color": GREY},
    {"indent": 1, "text": "ticker · date · source · texte",  "color": GREY, "size": 10},
    {"text": "Requête : similarité cosinus filtrée par ticker","color": WHITE, "bold": True},
    {"text": "Fallback automatique fetch live si index vide",  "color": ORANGE},
], Inches(0.35), Inches(1.42), Inches(4.3), line_h=Inches(0.3))

formula_box(s, Inches(0.35), Inches(3.7), Inches(4.3), Inches(0.45),
            "sim(q, d) = cos(E(q), E(d))  →  top-N docs renvoyés")

# Flèche
tb(s, "→", Inches(4.85), Inches(2.5), Inches(0.35), Inches(0.5),
   size=22, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# Étape G
rect(s, Inches(5.2), Inches(1.02), Inches(4.6), Inches(3.95), CARD)
tb(s, "② GENERATION — Groq LLaMA 3.3 70B", Inches(5.35), Inches(1.08), Inches(4.3), Inches(0.3),
   size=12, bold=True, color=ORANGE)
bullet_list(s, [
    {"text": "API Groq (inférence rapide, faible latence)",    "color": GREY},
    {"text": "Modèle : LLaMA 3.3 70B Versatile (défaut)",      "color": WHITE},
    {"text": "Temperature = 0.1 — quasi-déterministe",         "color": CYAN, "bold": True},
    {"text": "Prompt injecte :",                               "color": WHITE, "bold": True},
    {"indent": 1, "text": "ticker · date · régime GMM · headlines récupérés", "color": GREY, "size": 10},
    {"text": "Sortie JSON stricte :",                          "color": WHITE, "bold": True},
    {"indent": 1, "text": "sentiment  (bullish / bearish / neutral)", "color": GREEN, "size": 10},
    {"indent": 1, "text": "confidence  (0.0 – 1.0)",                 "color": GREEN, "size": 10},
    {"indent": 1, "text": "signal  (BUY / SELL / HOLD)",             "color": GREEN, "size": 10},
    {"indent": 1, "text": "reasoning  (texte explicatif)",           "color": GREY,  "size": 10},
], Inches(5.35), Inches(1.42), Inches(4.3), line_h=Inches(0.27))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Feature Matrix & Fusion
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "Feature Matrix & Fusion des Couches", "08",
       subtitle="Construction du vecteur 7D qui alimente les modèles XGBoost supervisés")
footer_bar(s)

# Vecteur central
rect(s, Inches(0.2), Inches(1.05), Inches(9.6), Inches(0.75), CARD)
tb(s, "Vecteur de Features — 7 dimensions hétérogènes",
   Inches(0.35), Inches(1.08), Inches(9.0), Inches(0.32),
   size=12, bold=True, color=CYAN)
tb(s, "feature_matrix_final.csv  —  90 lignes × 12 colonnes  —  30 jours × 3 tickers (AAPL · TSLA · NVDA)",
   Inches(0.35), Inches(1.4), Inches(9.0), Inches(0.3),
   size=9.5, color=GREY, italic=True)

features_vec = [
    ("RSI",          "14 périodes,\n0–100",              CYAN,   "Couche 1"),
    ("Volatilité",   "σ rolling 14j,\nretours quotidiens",ORANGE, "Couche 1"),
    ("MA Spread",    "(MA10-MA50)/MA50,\nnormalisé",      GREEN,  "Couche 1"),
    ("Régime GMM",   "0/1/2\nhaussier·calme·baissier",   PURPLE, "Couche 2A"),
    ("Sentiment",    "bullish/bearish/neutral\nencodé",   RED,    "Couche 2B"),
    ("Confiance",    "score LLM\n0.0 – 1.0",             ORANGE, "Couche 2B"),
    ("Signal RAG",   "BUY/SELL/HOLD\nencodé",             GREEN,  "Couche 2B"),
]
fx = Inches(0.2)
for fname, fdesc, fcol, flayer in features_vec:
    rect(s, fx, Inches(1.9), Inches(1.34), Inches(1.9), CARD)
    rect(s, fx, Inches(1.9), Inches(1.34), Pt(4), fcol)
    tb(s, fname, fx + Inches(0.05), Inches(1.94), Inches(1.24), Inches(0.4),
       size=10, bold=True, color=fcol, align=PP_ALIGN.CENTER)
    tb(s, fdesc, fx + Inches(0.05), Inches(2.34), Inches(1.24), Inches(0.55),
       size=8.5, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, flayer, fx + Inches(0.05), Inches(2.9), Inches(1.24), Inches(0.28),
       size=8, color=GREY, align=PP_ALIGN.CENTER, italic=True)
    fx += Inches(1.4)

# Flèche vers XGBoost
tb(s, "↓", Inches(4.6), Inches(3.9), Inches(0.8), Inches(0.4),
   size=22, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

rect(s, Inches(1.5), Inches(4.35), Inches(7.0), Inches(0.7), RED)
tb(s, "XGBoost Classifier  →  Prob. de hausse > 1%  |  XGBoost Regressor  →  Amplitude prédite",
   Inches(1.5), Inches(4.35), Inches(7.0), Inches(0.7),
   size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tb(s, "Cible classifieur : 1 si retour J+1 > 1%   |   Cible régresseur : retour brut J+1",
   Inches(0.2), Inches(5.1), Inches(9.6), Inches(0.28),
   size=9, color=GREY, italic=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — XGBoost Classifieur
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "XGBoost — Classifieur de Direction", "09",
       subtitle="Gradient Boosting supervisé : prédire si le titre va monter de plus de 1% le lendemain")
footer_bar(s)

# Principe Gradient Boosting
rect(s, Inches(0.2), Inches(1.02), Inches(5.6), Inches(2.2), CARD)
tb(s, "Principe — Gradient Boosting", Inches(0.35), Inches(1.08), Inches(5.3), Inches(0.3),
   size=12, bold=True, color=CYAN)
bullet_list(s, [
    {"text": "Ensemble d'arbres de décision entraînés séquentiellement", "color": WHITE, "bold": True},
    {"text": "Chaque arbre corrige les erreurs du précédent (résidus)", "color": GREY},
    {"text": "Régularisation L1/L2 intégrée — résistant au surapprentissage", "color": GREY},
    {"text": "Hyperparamètres choisis :", "color": WHITE, "bold": True},
    {"indent": 1, "text": "n_estimators=100  ·  max_depth=4  ·  learning_rate=0.1", "color": GREY, "size": 10},
], Inches(0.35), Inches(1.42), Inches(5.3), line_h=Inches(0.29))

tb(s, "Cible binaire :", Inches(0.35), Inches(2.65), Inches(5.3), Inches(0.28),
   size=11, bold=True, color=ORANGE)
formula_box(s, Inches(0.35), Inches(2.95), Inches(5.3), Inches(0.45),
            "y = 1  si  retour(t+1) > 1%   |   y = 0  sinon")

# Split temporel
rect(s, Inches(0.2), Inches(3.52), Inches(5.6), Inches(1.25), CARD)
tb(s, "Split temporel strict (no data leakage)", Inches(0.35), Inches(3.58), Inches(5.3), Inches(0.28),
   size=11, bold=True, color=RED)
# Barre de split
rect(s, Inches(0.35), Inches(3.93), Inches(3.7), Inches(0.35), GREEN)
rect(s, Inches(4.05), Inches(3.93), Inches(0.82), Inches(0.35), ORANGE)
rect(s, Inches(4.87), Inches(3.93), Inches(0.82), Inches(0.35), RED)
tb(s, "Train 70%", Inches(0.35), Inches(3.93), Inches(3.7), Inches(0.35),
   size=9, bold=True, color=DARK, align=PP_ALIGN.CENTER)
tb(s, "Val 15%", Inches(4.05), Inches(3.93), Inches(0.82), Inches(0.35),
   size=9, bold=True, color=DARK, align=PP_ALIGN.CENTER)
tb(s, "Test 15%", Inches(4.87), Inches(3.93), Inches(0.82), Inches(0.35),
   size=9, bold=True, color=DARK, align=PP_ALIGN.CENTER)
tb(s, "⚠ Jamais de split aléatoire sur données financières — look-ahead bias fatal",
   Inches(0.35), Inches(4.35), Inches(5.3), Inches(0.3),
   size=9, italic=True, color=RED)

# Métriques droite
metrics = [
    ("0.689", "AUC-ROC",   "Pouvoir discriminant\nau-dessus du hasard",   GREEN),
    ("0.491", "Recall",    "Sensibilité aux\nhausses > 1%",                ORANGE),
    ("0.060", "Precision", "Ratio vrais positifs\n(dataset déséquilibré)", RED),
    ("0.107", "F1-Score",  "Compromis\nprécision / rappel",                ORANGE),
]
mx = Inches(6.05)
my = Inches(1.02)
for i, (val, lbl, sub, col) in enumerate(metrics):
    cx = mx + (i % 2) * Inches(1.95)
    cy = my + (i // 2) * Inches(2.02)
    stat_box(s, cx, cy, Inches(1.72), Inches(1.78), val, lbl, col, sub)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — XGBoost Régresseur
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "XGBoost — Régresseur d'Amplitude", "10",
       subtitle="Prédire l'amplitude de la variation journalière pour calibrer la taille de position")
footer_bar(s)

rect(s, Inches(0.2), Inches(1.02), Inches(5.6), Inches(2.5), CARD)
tb(s, "Rôle du Régresseur", Inches(0.35), Inches(1.08), Inches(5.3), Inches(0.3),
   size=12, bold=True, color=ORANGE)
bullet_list(s, [
    {"text": "Complète le classifieur : la direction ne suffit pas", "color": WHITE, "bold": True},
    {"text": "Prédit le retour brut J+1 (valeur continue)", "color": GREY},
    {"text": "Utilisé par le backtester pour pondérer la position :", "color": WHITE},
    {"indent": 1, "text": "Signal BUY + amplitude > 1.5%  →  position boostée (×2.0 max)", "color": GREEN, "size": 10},
    {"indent": 1, "text": "Signal BUY standard  →  poids de base 1.2×", "color": GREY, "size": 10},
    {"indent": 1, "text": "Signal SELL  →  sortie de position (0×)", "color": RED, "size": 10},
    {"text": "Même architecture XGBoost, même feature matrix", "color": GREY},
    {"text": "Entraîné simultanément avec le classifieur", "color": GREY},
], Inches(0.35), Inches(1.42), Inches(5.3), line_h=Inches(0.27))

formula_box(s, Inches(0.35), Inches(3.62), Inches(5.3), Inches(0.5),
            "Cible : y = (Close(t+1) - Close(t)) / Close(t)  — retour logarithmique J+1")

# Métriques régresseur
rect(s, Inches(0.2), Inches(4.22), Inches(5.6), Inches(0.95), CARD2)
tb(s, "Métriques de régression (jeu de test) :", Inches(0.35), Inches(4.28), Inches(5.3), Inches(0.28),
   size=11, bold=True, color=ORANGE)
bullet_list(s, [
    {"text": "MAE  = 0.061   (Mean Absolute Error)",        "color": WHITE},
    {"text": "RMSE = 0.086   (Root Mean Squared Error)",    "color": WHITE},
    {"text": "Acc. directionnelle = 45.5%",                 "color": GREY},
], Inches(0.35), Inches(4.58), Inches(5.3), line_h=Inches(0.25))

# Droite — Sizing modèle
rect(s, Inches(5.9), Inches(1.02), Inches(3.9), Inches(4.15), CARD)
tb(s, "Position Sizing — Backtester", Inches(6.05), Inches(1.08), Inches(3.7), Inches(0.3),
   size=12, bold=True, color=GREEN)
sizing = [
    ("BUY + amplitude > 1.5%", "Poids ×2.0  (boost)", GREEN),
    ("BUY standard",            "Poids ×1.2  (base)",  CYAN),
    ("HOLD",                    "Poids ×1.0  (neutre)", GREY),
    ("SELL",                    "Poids ×0.0  (exit)",   RED),
]
sy = Inches(1.52)
for condition, weight, col in sizing:
    rect(s, Inches(6.05), sy, Inches(3.65), Inches(0.72), CARD2)
    tb(s, condition, Inches(6.12), sy + Inches(0.05), Inches(3.5), Inches(0.3),
       size=10, color=col, bold=True)
    tb(s, weight, Inches(6.12), sy + Inches(0.37), Inches(3.5), Inches(0.28),
       size=10, color=WHITE)
    sy += Inches(0.83)

tb(s, "Coûts de transaction : 0.1% commission + 0.05% slippage = 0.15%/trade",
   Inches(5.9), Inches(4.85), Inches(3.9), Inches(0.4),
   size=9, color=GREY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Backtesting & Résultats
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "Backtesting & Résultats", "11",
       subtitle="Simulation de portefeuille — Mars–Avril 2026 — Benchmark : S&P 500 (^GSPC)")
footer_bar(s)

big_stats = [
    ("2.625", "Sharpe Ratio",    "Excellent (> 1.2 = alpha)",       GREEN),
    ("60.0%", "Win Rate",        "60% de trades gagnants",          GREEN),
    ("+2.01%","Rendement Total", "Sur la période de test",          CYAN),
    ("-2.19%","Max Drawdown",    "Risque maximal observé",          ORANGE),
    ("48.4%", "Réduction Risque","vs S&P 500 sur même période",    GREEN),
]
sx = Inches(0.2)
for i, (val, lbl, sub, col) in enumerate(big_stats):
    stat_box(s, sx, Inches(1.05), Inches(1.86), Inches(1.75), val, lbl, col, sub)
    sx += Inches(1.93)

# Rappel métriques ML
rect(s, Inches(0.2), Inches(2.92), Inches(9.6), Inches(0.28), CARD)
tb(s, "Rappel métriques classifieur :", Inches(0.35), Inches(2.94), Inches(2.2), Inches(0.25),
   size=10, bold=True, color=CYAN)
tb(s, "AUC-ROC : 0.689   •   Recall : 0.491   •   F1 : 0.107   •   Precision : 0.060",
   Inches(2.6), Inches(2.94), Inches(7.0), Inches(0.25),
   size=10, color=WHITE)

# Modèle de backtest
rect(s, Inches(0.2), Inches(3.28), Inches(9.6), Inches(1.85), CARD)
tb(s, "Modèle de Backtest", Inches(0.35), Inches(3.34), Inches(9.2), Inches(0.3),
   size=12, bold=True, color=ORANGE)
bullet_list(s, [
    {"text": "Exécution : achat/vente au prix Open(T+1) — mesure P&L jusqu'au Close(T+1)", "color": WHITE, "bold": True},
    {"text": "Coûts : 0.1% commission + 0.05% slippage = 0.15% par trade",                "color": GREY},
    {"text": "Portefeuille initial : 10 000 $ répartis sur tous les tickers actifs",       "color": GREY},
    {"text": "Sharpe annualisé : √252 × (μ_retours / σ_retours)",                         "color": CYAN},
    {"text": "Statut «Alpha Generated» déclenché quand Sharpe > 1.2",                     "color": GREEN, "bold": True},
], Inches(0.35), Inches(3.68), Inches(9.2), line_h=Inches(0.27))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Pipeline Complet
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
header(s, "Pipeline Complet — De la donnée brute à la décision", "12")
footer_bar(s)

steps = [
    ("①\nDonnées",       "OHLCV + News\nyfinance",                  GREEN),
    ("②\nIndicateurs",   "RSI · Vol\nMA Spread",                    CYAN),
    ("③\nRégimes",       "GMM\nK=3",                                ORANGE),
    ("④\nAgent ReAct",   "3 iter.\nRAG + LLM",                      PURPLE),
    ("⑤\nFeature\nMatrix","7 features\n90 lignes",                  RED),
    ("⑥\nXGBoost\nTrain","Classif. +\nRégresseur",                  ORANGE),
    ("⑦\nPrédiction\nDouble","Direction\n+ Amplitude",              GREEN),
    ("⑧\nGUI\nElectron", "BUY/SELL/HOLD\nTemps réel",               CYAN),
]
sx = Inches(0.15)
for i, (title, desc, col) in enumerate(steps):
    rect(s, sx, Inches(1.02), Inches(1.2), Inches(2.1), col)
    tb(s, title, sx, Inches(1.02), Inches(1.2), Inches(1.2),
       size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    tb(s, desc, sx, Inches(2.15), Inches(1.2), Inches(0.95),
       size=8.5, color=DARK, align=PP_ALIGN.CENTER)
    if i < 7:
        tb(s, "→", sx + Inches(1.21), Inches(1.78), Inches(0.18), Inches(0.4),
           size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sx += Inches(1.24)

details = [
    ("Entrée utilisateur",
     "Ticker + Date  →  inference.py déclenche l'agent ReAct (3 itérations max)", GREY),
    ("Agent ReAct",
     "get_regime()  →  search_news_db()  →  Groq LLaMA → JSON", CYAN),
    ("Feature Matrix",
     "RSI + Vol + MA + regime_id + sentiment + confidence + rag_signal  →  vecteur 7D", ORANGE),
    ("XGBoost Output",
     "Direction (HAUSSE/BAISSE) + Probabilité % + Amplitude prédite ±xx%", GREEN),
    ("GUI Electron",
     "Affichage BUY/SELL/HOLD + Agent Thought Stream en temps réel via FastAPI", WHITE),
]
dy = Inches(3.25)
for dlabel, dtext, dcol in details:
    rect(s, Inches(0.2), dy, Inches(9.6), Inches(0.35), CARD2)
    tb(s, dlabel, Inches(0.3), dy + Inches(0.03), Inches(2.1), Inches(0.29),
       size=9, bold=True, color=dcol)
    tb(s, dtext, Inches(2.45), dy + Inches(0.03), Inches(7.3), Inches(0.29),
       size=9, color=WHITE)
    dy += Inches(0.41)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conclusion
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
rect(s, 0, 0, SW, Inches(0.9), CARD)
rect(s, 0, Inches(0.9), SW, Pt(3), CYAN)
tb(s, "Conclusion", Inches(0.3), Inches(0.06), Inches(9.4), Inches(0.72),
   size=26, bold=True, color=CYAN)
footer_bar(s)

tb(s, "Sentirade-RAG",
   Inches(0.3), Inches(1.02), Inches(9.4), Inches(0.45),
   size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Un pipeline ML complet, multi-modal, de bout en bout",
   Inches(0.3), Inches(1.47), Inches(9.4), Inches(0.35),
   size=14, italic=True, color=GREY, align=PP_ALIGN.CENTER)

checkpoints = [
    ("Agent ReAct RAG justifié par decision_log.csv (450+ lignes de log autonomes)",           CYAN),
    ("Split temporel strict — jamais de data leakage sur données financières",                  RED),
    ("Double prédiction : direction (XGBoost Classifier) + amplitude (XGBoost Regressor)",     ORANGE),
    ("GMM non supervisé pour contextualiser le marché avant l'appel LLM",                      GREEN),
    ("Interface Electron + FastAPI + yfinance temps réel — démo live possible",                 WHITE),
    ("AUC-ROC 0.689 · Sharpe Ratio 2.625 · Win Rate 60% · Réduction risque 48.4%",            GREEN),
]
cy = Inches(1.98)
for text, col in checkpoints:
    rect(s, Inches(0.3), cy, Inches(0.32), Inches(0.32), col)
    tb(s, "✓", Inches(0.3), cy, Inches(0.32), Inches(0.32),
       size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    tb(s, text, Inches(0.7), cy, Inches(9.0), Inches(0.32),
       size=11, color=col, bold=(col != GREY))
    cy += Inches(0.42)

rect(s, 0, Inches(5.15), SW, Inches(0.2), DARK)
tb(s, "ISITD 2024 — S4 Machine Learning   |   github.com/wia-mej/sentirade-rag1",
   Inches(0.3), Inches(5.16), Inches(9.4), Inches(0.19),
   size=8, color=GREY, align=PP_ALIGN.CENTER)

# ── Sauvegarde ─────────────────────────────────────────────────────────────────
out = "/home/zaid/Documents/ML/sentirade-rag1/Sentirade_RAG_Presentation_FR.pptx"
prs.save(out)
print(f"✅  Sauvegardé : {out}")
print(f"    Slides : {len(prs.slides)}")
